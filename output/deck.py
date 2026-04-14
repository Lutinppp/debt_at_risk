"""
PowerPoint deck builder for the EU G4 Debt-at-Risk board presentation.

Produces a 5-slide deck with the institutional navy (#1B2A4A) and
gold (#C8A951) colour palette using python-pptx.

Slides
------
1. Global context — G4 in world debt distribution
2. Fan charts for FR, DE, IT, ES
3. Drivers waterfall per country
4. Fiscal crisis probability scores
5. Policy implications (bullet points, editable)

Output: ``output/eu_g4_debt_at_risk.pptx``
"""

from __future__ import annotations

import logging
from io import BytesIO
from pathlib import Path
from typing import Any

import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)

OUTPUT_DIR = Path(__file__).parent
PPTX_FILE = OUTPUT_DIR / "eu_g4_debt_at_risk.pptx"
CHARTS_DIR = OUTPUT_DIR / "charts"

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY_RGB = RGBColor(0x1B, 0x2A, 0x4A)
GOLD_RGB = RGBColor(0xC8, 0xA9, 0x51)
WHITE_RGB = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY_RGB = RGBColor(0xEC, 0xF0, 0xF1)


def _slide_width() -> float:
    return 10.0  # inches


def _slide_height() -> float:
    return 7.5  # inches


def _add_slide(prs: Presentation, layout_idx: int = 6) -> Any:
    """Add a blank slide (layout 6 = blank in most themes)."""
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)


def _set_bg(slide: Any, color: RGBColor = NAVY_RGB) -> None:
    """Set slide background to solid colour."""
    from pptx.oxml.ns import qn  # noqa: PLC0415
    from lxml import etree  # noqa: PLC0415

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(
    slide: Any,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int = 14,
    bold: bool = False,
    color: RGBColor = WHITE_RGB,
    align: Any = PP_ALIGN.LEFT,
    wrap: bool = True,
) -> Any:
    from pptx.util import Inches, Pt  # noqa: PLC0415

    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def _add_image_from_fig(
    slide: Any,
    fig: plt.Figure,
    left: float,
    top: float,
    width: float,
    height: float,
) -> None:
    """Save a matplotlib figure to a BytesIO stream and insert into slide."""
    stream = BytesIO()
    fig.savefig(stream, format="png", dpi=120, bbox_inches="tight", facecolor="white")
    stream.seek(0)
    slide.shapes.add_picture(
        stream,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    plt.close(fig)


def _add_image_from_path(
    slide: Any,
    path: Path,
    left: float,
    top: float,
    width: float,
    height: float,
) -> None:
    if not path.exists():
        logger.warning("Chart file not found: %s", path)
        return
    slide.shapes.add_picture(
        str(path),
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )


def _add_divider(
    slide: Any,
    top: float,
    color: RGBColor = GOLD_RGB,
    width: float = 9.0,
    left: float = 0.5,
    height: float = 0.04,
) -> None:
    from pptx.util import Inches, Pt  # noqa: PLC0415

    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------


def _slide1_global_context(prs: Presentation, panel: Any | None = None) -> None:
    """Slide 1: Global context — G4 in world debt distribution."""
    slide = _add_slide(prs)
    _set_bg(slide, NAVY_RGB)

    _add_textbox(
        slide,
        "EU G4 Debt-at-Risk: Global Context",
        left=0.5,
        top=0.3,
        width=9.0,
        height=0.6,
        font_size=24,
        bold=True,
        color=GOLD_RGB,
        align=PP_ALIGN.LEFT,
    )
    _add_divider(slide, top=1.05)

    # Body text
    body = (
        "• Global government debt has risen to historic highs, averaging ~93% of GDP in 2024.\n"
        "• EU G4 (France, Germany, Italy, Spain) account for ~60% of EU GDP.\n"
        "• Debt trajectories diverge sharply: Italy and France face elevated risk;\n"
        "  Germany benefits from fiscal headroom; Spain shows moderate risk.\n"
        "• This analysis applies the IMF Debt-at-Risk (DaR) framework\n"
        "  (Furceri, Giannone, Kisat, Lam, Li — WP/25/86) to quantify\n"
        "  the distribution of future debt paths under alternative scenarios.\n\n"
        "Source: IMF WEO (April 2025), ECB SDW, IMF FSI, WUI."
    )
    _add_textbox(
        slide,
        body,
        left=0.5,
        top=1.2,
        width=9.0,
        height=5.5,
        font_size=13,
        color=WHITE_RGB,
    )

    # Watermark reference
    _add_textbox(
        slide,
        "IMF WP/25/86 — Furceri, Giannone, Kisat, Lam, Li (May 2025)",
        left=0.5,
        top=7.0,
        width=9.0,
        height=0.3,
        font_size=8,
        color=GOLD_RGB,
        align=PP_ALIGN.RIGHT,
    )


def _slide2_fan_charts(prs: Presentation) -> None:
    """Slide 2: Fan charts for FR, DE, IT, ES."""
    slide = _add_slide(prs)
    _set_bg(slide, NAVY_RGB)

    _add_textbox(
        slide,
        "Government Debt Fan Charts — EU G4 (2027 Horizon)",
        left=0.5,
        top=0.2,
        width=9.0,
        height=0.5,
        font_size=20,
        bold=True,
        color=GOLD_RGB,
        align=PP_ALIGN.LEFT,
    )
    _add_divider(slide, top=0.85)

    chart_path = CHARTS_DIR / "fan_charts.png"
    _add_image_from_path(slide, chart_path, left=0.4, top=1.0, width=9.2, height=6.0)

    _add_textbox(
        slide,
        "Shaded band: P5–P95 distribution of debt/GDP.  Dashed line: P50 median.  Dotted red: P95 (Debt-at-Risk).",
        left=0.5,
        top=7.1,
        width=9.0,
        height=0.3,
        font_size=8,
        color=LIGHT_GREY_RGB,
    )


def _slide3_waterfall(prs: Presentation) -> None:
    """Slide 3: Drivers waterfall per country."""
    slide = _add_slide(prs)
    _set_bg(slide, NAVY_RGB)

    _add_textbox(
        slide,
        "Upside Debt Risk Decomposition by Driver",
        left=0.5,
        top=0.2,
        width=9.0,
        height=0.5,
        font_size=20,
        bold=True,
        color=GOLD_RGB,
    )
    _add_divider(slide, top=0.85)

    # 2×2 grid of waterfall charts
    positions = [
        ("FRA", 0.3, 1.0),
        ("DEU", 5.1, 1.0),
        ("ITA", 0.3, 4.1),
        ("ESP", 5.1, 4.1),
    ]
    for iso, left, top in positions:
        chart_path = CHARTS_DIR / f"waterfall_{iso}.png"
        _add_image_from_path(slide, chart_path, left=left, top=top, width=4.5, height=3.0)

    _add_textbox(
        slide,
        "Bars show contribution (pp of GDP) from each conditioning variable to P95−P50 upside risk at h=3y.",
        left=0.5,
        top=7.1,
        width=9.0,
        height=0.3,
        font_size=8,
        color=LIGHT_GREY_RGB,
    )


def _slide4_crisis_signals(prs: Presentation) -> None:
    """Slide 4: Fiscal crisis probability scores."""
    slide = _add_slide(prs)
    _set_bg(slide, NAVY_RGB)

    _add_textbox(
        slide,
        "Fiscal Crisis Early-Warning Signals — 2025–2026",
        left=0.5,
        top=0.2,
        width=9.0,
        height=0.5,
        font_size=20,
        bold=True,
        color=GOLD_RGB,
    )
    _add_divider(slide, top=0.85)

    chart_path = CHARTS_DIR / "crisis_signals.png"
    _add_image_from_path(slide, chart_path, left=1.5, top=1.1, width=7.0, height=4.5)

    _add_textbox(
        slide,
        (
            "Fiscal crisis probability estimated via panel logit:\n"
            "  P(crisis_{t+1}) = Λ(α + β · Upside_{t})  [Upside = P95 − P50]\n"
            "Crisis events: IMF / Laeven-Valencia (2020) database.\n"
            "Red dashed line = 10% early-warning threshold."
        ),
        left=0.5,
        top=5.8,
        width=9.0,
        height=1.5,
        font_size=9,
        color=LIGHT_GREY_RGB,
    )


def _slide5_policy(prs: Presentation) -> None:
    """Slide 5: Policy implications (bullet points, editable)."""
    slide = _add_slide(prs)
    _set_bg(slide, NAVY_RGB)

    _add_textbox(
        slide,
        "Policy Implications",
        left=0.5,
        top=0.2,
        width=9.0,
        height=0.5,
        font_size=22,
        bold=True,
        color=GOLD_RGB,
    )
    _add_divider(slide, top=0.85)

    bullets = [
        "1.  Italy and France face elevated Debt-at-Risk (P95 > 140% and > 125% GDP respectively by 2027).",
        "",
        "2.  Financial stress and sovereign spreads are the dominant drivers of upside risk for Italy and Spain.",
        "",
        "3.  Germany retains significant fiscal space; DaR remains below 75% GDP at the 3-year horizon.",
        "",
        "4.  Fiscal crisis probability signals warrant enhanced monitoring for Italy (>15%) and Spain (>10%).",
        "",
        "5.  Policy priority: credible medium-term fiscal consolidation to reduce Upside risk and lower DaR.",
        "",
        "6.  International coordination on financial stability buffers can reduce correlated tail risks across G4.",
        "",
        "[EDITABLE — insert updated forecasts, policy actions, or board commentary here]",
    ]
    _add_textbox(
        slide,
        "\n".join(bullets),
        left=0.5,
        top=1.1,
        width=9.0,
        height=6.0,
        font_size=12,
        color=WHITE_RGB,
    )

    _add_textbox(
        slide,
        "Methodology: IMF WP/25/86 — Furceri, Giannone, Kisat, Lam, Li (May 2025)",
        left=0.5,
        top=7.1,
        width=9.0,
        height=0.3,
        font_size=8,
        color=GOLD_RGB,
        align=PP_ALIGN.RIGHT,
    )


# ---------------------------------------------------------------------------
# Main deck builder
# ---------------------------------------------------------------------------


def build_deck(
    panel: Any = None,
    save: bool = True,
    output_path: Path | None = None,
) -> Presentation:
    """Build the 5-slide PowerPoint deck.

    Parameters
    ----------
    panel : pd.DataFrame, optional
        Historical panel data (used for slide 1 context).
    save : bool
        Write the file to ``output/eu_g4_debt_at_risk.pptx`` when *True*.
    output_path : Path, optional
        Override the default output path.

    Returns
    -------
    Presentation
        The python-pptx Presentation object.
    """
    prs = Presentation()
    prs.slide_width = Inches(_slide_width())
    prs.slide_height = Inches(_slide_height())

    logger.info("Building slide 1: global context …")
    _slide1_global_context(prs, panel=panel)

    logger.info("Building slide 2: fan charts …")
    _slide2_fan_charts(prs)

    logger.info("Building slide 3: waterfall drivers …")
    _slide3_waterfall(prs)

    logger.info("Building slide 4: crisis signals …")
    _slide4_crisis_signals(prs)

    logger.info("Building slide 5: policy implications …")
    _slide5_policy(prs)

    if save:
        dest = output_path or PPTX_FILE
        dest.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(dest))
        logger.info("Saved deck to %s", dest)

    return prs


if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO, format="%(levelname)s %(message)s")
    build_deck()
